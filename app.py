import streamlit as st
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import io
import re

# ReportLab Imports
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# 1. Page Configuration
st.set_page_config(
    page_title="Laaura's Resume & Document Summarizer",
    page_icon="🩰",
    layout="wide"
)

# 2. Themes Definition
THEMES = {
    "🩰 Coquette Pink": {
        "gradient": "linear-gradient(135deg, #FFE4E1 0%, #FFC0CB 50%, #E6E6FA 100%)",
        "paper_gradient": "linear-gradient(135deg, rgba(255, 255, 255, 0.95) 0%, rgba(255, 228, 225, 0.90) 50%, rgba(255, 192, 203, 0.85) 100%)",
        "card_bg": "rgba(255, 255, 255, 0.75)",
        "text": "#5C2C3B",
        "accent": "#E899AC",
        "emojis": "🩰 🎀 🕯️ 🦢 🎀",
        "pdf_bg": colors.HexColor("#FFF0F5"),
        "pdf_border": colors.HexColor("#FFB6C1"),
        "pdf_header": colors.HexColor("#D87093"),
        "line_color": "#FFC0CB"
    },
    "🍷 Maroon Strawberry": {
        "gradient": "linear-gradient(135deg, #3B0008 0%, #5E0015 50%, #800020 100%)",
        "paper_gradient": "linear-gradient(135deg, rgba(255, 255, 255, 0.95) 0%, rgba(250, 235, 238, 0.92) 50%, rgba(244, 204, 211, 0.88) 100%)",
        "card_bg": "rgba(255, 255, 255, 0.15)",
        "text": "#FFFFFF",
        "accent": "#C72C41",
        "emojis": "🍷 🍓 🥀 🕯️ 🍒",
        "pdf_bg": colors.HexColor("#FAEBEE"),
        "pdf_border": colors.HexColor("#C72C41"),
        "pdf_header": colors.HexColor("#800020"),
        "line_color": "#E8B4B8"
    },
    "☁️ Langit Cerah": {
        "gradient": "linear-gradient(135deg, #E0F7FA 0%, #B3E5FC 50%, #E1BEE7 100%)",
        "paper_gradient": "linear-gradient(135deg, rgba(255, 255, 255, 0.95) 0%, rgba(224, 247, 250, 0.90) 50%, rgba(179, 229, 252, 0.85) 100%)",
        "card_bg": "rgba(255, 255, 255, 0.75)",
        "text": "#1F3A52",
        "accent": "#81D4FA",
        "emojis": "☁️ 🌙 🕊️ 💫 🌤️",
        "pdf_bg": colors.HexColor("#F0F8FF"),
        "pdf_border": colors.HexColor("#B0E0E6"),
        "pdf_header": colors.HexColor("#0288D1"),
        "line_color": "#B0E0E6"
    },
    "🍵 Matcha Soft": {
        "gradient": "linear-gradient(135deg, #F1F8E9 0%, #DCEDC8 50%, #C8E6C9 100%)",
        "paper_gradient": "linear-gradient(135deg, rgba(255, 255, 255, 0.95) 0%, rgba(241, 248, 233, 0.90) 50%, rgba(220, 237, 200, 0.85) 100%)",
        "card_bg": "rgba(255, 255, 255, 0.75)",
        "text": "#2D4A27",
        "accent": "#AED581",
        "emojis": "🍵 🍃 🌱 🫖 🧁",
        "pdf_bg": colors.HexColor("#F4F9F4"),
        "pdf_border": colors.HexColor("#C8E6C9"),
        "pdf_header": colors.HexColor("#558B2F"),
        "line_color": "#C8E6C9"
    },
    "🍂 Earth Warm": {
        "gradient": "linear-gradient(135deg, #D7C4B7 0%, #AF9483 50%, #8C6D58 100%)",
        "paper_gradient": "linear-gradient(135deg, rgba(255, 255, 255, 0.95) 0%, rgba(245, 235, 226, 0.90) 50%, rgba(230, 215, 201, 0.85) 100%)",
        "card_bg": "rgba(255, 255, 255, 0.75)",
        "text": "#3E2723",
        "accent": "#8D6E63",
        "emojis": "🍂 ☕ 🧸 🪵 🌾",
        "pdf_bg": colors.HexColor("#FAF0E6"),
        "pdf_border": colors.HexColor("#D7B596"),
        "pdf_header": colors.HexColor("#5D4037"),
        "line_color": "#D7B596"
    }
}

FONTS = {
    "Poppins (Modern Clean)": {"rl": "Helvetica", "rl_bold": "Helvetica-Bold", "css": "'Poppins', sans-serif"},
    "Times New Roman / Series": {"rl": "Times-Roman", "rl_bold": "Times-Bold", "css": "'Times New Roman', Times, serif"},
    "Classic Serif (Times)": {"rl": "Times-Roman", "rl_bold": "Times-Bold", "css": "'Times New Roman', serif"},
    "Modern Sans (Helvetica)": {"rl": "Helvetica", "rl_bold": "Helvetica-Bold", "css": "'Helvetica Neue', Helvetica, Arial, sans-serif"},
    "Typewriter Monospace (Courier)": {"rl": "Courier", "rl_bold": "Courier-Bold", "css": "'Courier New', Courier, monospace"}
}

FONT_SIZES = {
    "10 pt": {"css": "10pt", "rl": 10, "leading": 14},
    "11 pt": {"css": "11pt", "rl": 11, "leading": 15},
    "12 pt": {"css": "12pt", "rl": 12, "leading": 16},
    "14 pt": {"css": "14pt", "rl": 14, "leading": 18},
    "16 pt": {"css": "16pt", "rl": 16, "leading": 20},
    "18 pt": {"css": "18pt", "rl": 18, "leading": 22}
}

# 3. Sidebar Customization
st.sidebar.subheader("🎨 Theme")
selected_theme_name = st.sidebar.selectbox("Pilih Tema Color Palette:", list(THEMES.keys()))
theme = THEMES[selected_theme_name]
emojis = theme['emojis'].split()

st.sidebar.subheader("📄 Kustom Kertas")
paper_style = st.sidebar.selectbox(
    "Pilih Kertas:",
    ["Polos (Clean Blank)", "Buku Tulis (Ruled Lines)", "Kotak-Kotak (Grid)", "Bintik-Bintik (Dotted)"]
)

font_style_name = st.sidebar.selectbox("Pilih Gaya Font:", list(FONTS.keys()))
selected_font = FONTS[font_style_name]

font_size_name = st.sidebar.selectbox("Ukuran Font (pt):", list(FONT_SIZES.keys()), index=2)
selected_size = FONT_SIZES[font_size_name]

st.sidebar.subheader("📝 Rangkuman")
summary_format = st.sidebar.selectbox(
    "Format Rangkuman:",
    ["Struktur Rapi AI", "Bullet Points (Poin-Poin Pilihan)", "Paragraf Eksekutif"]
)
summary_length = st.sidebar.radio("Panjang Rangkuman:", ["Singkat", "Sedang", "Lengkap"])

# Inject CSS Dynamic Styling
paper_pattern_css = ""
if paper_style == "Buku Tulis (Ruled Lines)":
    paper_pattern_css = f", repeating-linear-gradient(transparent, transparent 27px, {theme['line_color']} 28px)"
elif paper_style == "Kotak-Kotak (Grid)":
    paper_pattern_css = f", linear-gradient({theme['line_color']} 1px, transparent 1px), linear-gradient(90deg, {theme['line_color']} 1px, transparent 1px)"
elif paper_style == "Bintik-Bintik (Dotted)":
    paper_pattern_css = f", radial-gradient({theme['line_color']} 1.5px, transparent 1.5px)"

st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700&display=swap');

    .stApp, [data-testid="stSidebar"] {{
        background: {theme['gradient']} !important;
        background-attachment: fixed !important;
    }}

    .stApp p, .stApp h1, .stApp h2, .stApp h3, .stApp h4, .stApp label {{
        color: {theme['text']};
        font-family: {selected_font['css']};
    }}

    [data-testid="stSidebarContent"] {{
        padding: 15px !important;
    }}

    [data-testid="stSidebar"] div.stSelectbox, 
    [data-testid="stSidebar"] div.stRadio {{
        background-color: rgba(255, 255, 255, 0.95) !important;
        border-radius: 12px !important;
        padding: 12px 15px !important;
        margin-bottom: 12px !important;
        box-shadow: 0 4px 10px rgba(0,0,0,0.1) !important;
    }}

    [data-testid="stSidebar"] *, 
    [data-testid="stSidebar"] p, 
    [data-testid="stSidebar"] span, 
    [data-testid="stSidebar"] label, 
    [data-testid="stSidebar"] h1, 
    [data-testid="stSidebar"] h2, 
    [data-testid="stSidebar"] h3, 
    [data-testid="stSidebar"] div,
    [data-testid="stSidebar"] option {{
        color: #000000 !important;
        font-weight: 500;
    }}

    [data-testid="stSidebar"] h3 {{
        background-color: rgba(255, 255, 255, 0.95) !important;
        padding: 8px 12px !important;
        border-radius: 8px !important;
        font-weight: 700 !important;
        font-size: 1.05rem !important;
        margin-top: 10px !important;
        margin-bottom: 10px !important;
    }}

    .main-title {{
        text-align: center;
        color: {theme['text']} !important;
        font-weight: 700;
        font-size: 2.5rem;
        margin-bottom: 0px;
        font-family: {selected_font['css']};
    }}
    
    .watermark-tag {{
        text-align: center;
        font-size: 1.1rem;
        color: {theme['text']} !important;
        font-weight: 600;
        letter-spacing: 1px;
        margin-top: 5px;
        margin-bottom: 25px;
        opacity: 0.9;
        font-family: {selected_font['css']};
    }}

    div[data-testid="stFileUploader"] {{
        background-color: rgba(255, 255, 255, 0.95) !important;
        border-radius: 15px !important;
        padding: 15px !important;
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }}

    div[data-testid="stFileUploader"] label,
    div[data-testid="stFileUploader"] label p {{
        color: #000000 !important;
        font-weight: bold !important;
    }}

    [data-testid="stFileUploaderDropzone"] {{
        background-color: #FFFFFF !important;
        border: 2px dashed {theme['accent']} !important;
        border-radius: 12px !important;
        padding: 15px !important;
    }}

    [data-testid="stFileUploaderDropzone"] * {{
        color: #000000 !important;
    }}

    /* Layout Rangkuman Rapi ala GPT/Gemini */
    .preview-paper {{
        background: {theme['paper_gradient']} {paper_pattern_css};
        background-size: cover, 20px 20px, 20px 20px;
        border: 2px solid {theme['line_color']};
        border-radius: 16px;
        padding: 35px;
        box-shadow: 0 10px 25px rgba(0,0,0,0.08);
        color: #2D3748 !important;
        font-family: {selected_font['css']};
        font-size: {selected_size['css']};
        line-height: 1.8;
    }}
    
    .preview-paper h3 {{
        color: {theme['text']} !important;
        font-size: 1.25em !important;
        margin-top: 20px !important;
        margin-bottom: 10px !important;
        border-bottom: 2px solid {theme['line_color']};
        padding-bottom: 5px;
        font-weight: 700 !important;
    }}

    .preview-paper ul {{
        margin-top: 5px;
        margin-bottom: 15px;
        padding-left: 20px;
    }}

    .preview-paper li {{
        margin-bottom: 8px;
        line-height: 1.6;
    }}

    .stButton>button {{
        background-color: {theme['accent']} !important;
        color: white !important;
        border-radius: 20px !important;
        border: none !important;
        font-weight: bold !important;
        padding: 0.6rem 1.8rem !important;
        font-size: 1rem !important;
        box-shadow: 0 4px 10px rgba(0,0,0,0.08);
    }}
    </style>
""", unsafe_allow_html=True)

# 4. Header Section
st.markdown(f"<h1 class='main-title'>{emojis[0]} Laaura's Resume {emojis[1]}</h1>", unsafe_allow_html=True)
st.markdown(f"<p class='watermark-tag'>✨ Smart Document & Resume Summarizer ✨</p>", unsafe_allow_html=True)

# 5. Text Extraction Functions
def extract_text_from_pdf(uploaded_file):
    text = ""
    doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    for page in doc:
        text += page.get_text() + "\n"
    return text

def extract_text_from_docx(uploaded_file):
    doc = Document(uploaded_file)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])

def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

# 6. Smart AI-Style Summarizer (Rapi ala GPT / Gemini)
def format_sentence_highlight(sentence):
    words = sentence.split()
    if len(words) > 3:
        # Menjadikan 2-4 kata pertama bold sebagai poin kunci (style GPT/Gemini)
        head = " ".join(words[:3])
        tail = " ".join(words[3:])
        return f"<b>{head}</b> {tail}"
    return sentence

def generate_ai_style_summary(text, length_option, format_option):
    # Membersihkan kalimat
    raw_sentences = [s.strip() for s in re.split(r'[.\n]+', text) if len(s.strip()) > 12]
    if not raw_sentences:
        return "<i>Tidak ada teks yang cukup untuk dirangkum dari dokumen ini.</i>"

    if length_option == "Singkat":
        limit = max(3, len(raw_sentences) // 4)
    elif length_option == "Sedang":
        limit = max(6, len(raw_sentences) // 2)
    else:
        limit = max(9, int(len(raw_sentences) * 0.75))
        
    sentences = raw_sentences[:limit]

    if format_option == "Bullet Points (Poin-Poin Pilihan)":
        bullets = "".join([f"<li>{format_sentence_highlight(s)}.</li>" for s in sentences])
        return f"<ul>{bullets}</ul>"
        
    elif format_option == "Paragraf Eksekutif":
        paras = ". ".join(sentences) + "."
        return f"<p>{paras}</p>"
        
    else:  # Format "Struktur Rapi AI (GPT & Gemini Style)"
        n = len(sentences)
        sec1_end = max(1, n // 3)
        sec2_end = max(sec1_end + 1, (2 * n) // 3)

        overview_sents = sentences[:sec1_end]
        points_sents = sentences[sec1_end:sec2_end]
        action_sents = sentences[sec2_end:]

        html_out = "<h3>📌 Ringkasan Eksekutif & Gambaran Umum</h3>"
        html_out += f"<p>{'. '.join(overview_sents)}.</p>"

        html_out += "<h3>🔑 Poin-Poin Kunci & Highlight Dokumen</h3><ul>"
        for s in points_sents:
            html_out += f"<li>{format_sentence_highlight(s)}.</li>"
        html_out += "</ul>"

        if action_sents:
            html_out += "<h3>💡 Insight & Kesimpulan Utama</h3><ul>"
            for s in action_sents:
                html_out += f"<li>{format_sentence_highlight(s)}.</li>"
            html_out += "</ul>"

        return html_out

# 7. ReportLab PDF Engine (Mendukung HTML Rapi GPT Style)
def create_custom_pdf(summary_html, custom_uploaded_images, theme_info, font_info, size_info, pattern_name):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, 
        pagesize=letter, 
        rightMargin=36, 
        leftMargin=36, 
        topMargin=36, 
        bottomMargin=36
    )
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'HeaderTitleCustom',
        fontName=font_info['rl_bold'],
        fontSize=18,
        textColor=theme_info['pdf_header'],
        alignment=1,
        spaceAfter=4
    )
    
    sub_style = ParagraphStyle(
        'HeaderSubCustom',
        fontName=font_info['rl'],
        fontSize=10,
        textColor=colors.gray,
        alignment=1
    )

    h3_style = ParagraphStyle(
        'H3StylePDF',
        fontName=font_info['rl_bold'],
        fontSize=size_info['rl'] + 2,
        leading=size_info['leading'] + 2,
        textColor=theme_info['pdf_header'],
        spaceBefore=10,
        spaceAfter=6
    )

    body_style = ParagraphStyle(
        'BodyTextCustomPaper',
        fontName=font_info['rl'],
        fontSize=size_info['rl'],
        leading=size_info['leading'],
        textColor=colors.HexColor("#2D3748"),
        spaceAfter=6
    )

    story = []

    # Header Card
    header_data = [[
        Paragraph(f"<b>{theme_info['emojis']} Laaura's Resume {theme_info['emojis']}</b>", title_style),
    ], [
        Paragraph("<b>✨ Smart Document & Resume Summarizer ✨</b>", sub_style)
    ]]
    
    header_table = Table(header_data, colWidths=[520])
    header_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.white),
        ('BOX', (0,0), (-1,-1), 1.5, theme_info['pdf_border']),
        ('PADDING', (0,0), (-1,-1), 10),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
    ]))
    
    story.append(header_table)
    story.append(Spacer(1, 14))

    # Parse HTML ringkasan ke bentuk komponen ReportLab
    # Konversi tag HTML dasar ke elemen ReportLab
    blocks = re.split(r'(<h3>.*?</h3>)', summary_html)
    pdf_flowables = []

    for block in blocks:
        if not block.strip():
            continue
        if block.startswith('<h3>'):
            clean_title = block.replace('<h3>', '').replace('</h3>', '')
            pdf_flowables.append(Paragraph(clean_title, h3_style))
        elif '<ul>' in block:
            items = re.findall(r'<li>(.*?)</li>', block)
            for item in items:
                pdf_flowables.append(Paragraph(f"• {item}", body_style))
        else:
            clean_p = block.replace('<p>', '').replace('</p>', '').strip()
            if clean_p:
                pdf_flowables.append(Paragraph(clean_p, body_style))

    # Wrap ringkasan dalam box tabel rapi
    summary_table = Table([[pdf_flowables]], colWidths=[520])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.white),
        ('LINELEFT', (0,0), (-1,-1), 4, theme_info['pdf_header']),
        ('BOX', (0,0), (-1,-1), 1, theme_info['pdf_border']),
        ('PADDING', (0,0), (-1,-1), 14),
    ]))
    
    story.append(summary_table)

    # Manual Uploaded Images Section
    if custom_uploaded_images:
        story.append(Spacer(1, 15))
        story.append(Paragraph(f"<b>{emojis[2]} Lampiran Gambar Tambahan:</b>", ParagraphStyle('ImgHeaderCustom', fontName=font_info['rl_bold'], fontSize=11, textColor=theme_info['pdf_header'])))
        story.append(Spacer(1, 8))
        
        for img_file in custom_uploaded_images:
            img = Image.open(img_file)
            img_buffer = io.BytesIO()
            img.convert('RGB').save(img_buffer, format='JPEG')
            img_buffer.seek(0)
            rl_img = RLImage(img_buffer, width=280, height=170)
            
            img_table = Table([[rl_img]], colWidths=[520])
            img_table.setStyle(TableStyle([
                ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ('PADDING', (0,0), (-1,-1), 6)
            ]))
            story.append(img_table)
            story.append(Spacer(1, 10))

    def draw_background_and_pattern(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(theme_info['pdf_bg'])
        canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=1, stroke=0)
        
        line_col = theme_info['pdf_border']
        canvas.setStrokeColor(line_col)
        canvas.setFillColor(line_col)
        canvas.setLineWidth(0.5)
        
        width, height = doc.pagesize
        
        if pattern_name == "Buku Tulis (Ruled Lines)":
            y = 20
            while y < height:
                canvas.line(0, y, width, y)
                y += 24
        elif pattern_name == "Kotak-Kotak (Grid)":
            x = 0
            while x < width:
                canvas.line(x, 0, x, height)
                x += 20
            y = 0
            while y < height:
                canvas.line(0, y, width, y)
                y += 20
        elif pattern_name == "Bintik-Bintik (Dotted)":
            x = 10
            while x < width:
                y = 10
                while y < height:
                    canvas.circle(x, y, 1, fill=1, stroke=0)
                    y += 20
                x += 20
                
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_background_and_pattern, onLaterPages=draw_background_and_pattern)
    buffer.seek(0)
    return buffer

# 8. Interactive App Workflow
col1, col2 = st.columns([2, 1])

with col1:
    uploaded_file = st.file_uploader(
        f"{emojis[3]} Unggah dokumen kamu (PDF, PPTX, atau DOCX):", 
        type=["pdf", "pptx", "docx"]
    )

with col2:
    uploaded_images = st.file_uploader(
        f"{emojis[2]} Upload Gambar Sendiri (Optional):", 
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True
    )

if uploaded_file is not None:
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    if st.button("✨ Buat Rangkuman Custom"):
        with st.spinner(f"Membaca dan merangkum dokumen kamu... {emojis[4]}"):
            if file_type == "pdf":
                raw_text = extract_text_from_pdf(uploaded_file)
            elif file_type == "docx":
                raw_text = extract_text_from_docx(uploaded_file)
            elif file_type == "pptx":
                raw_text = extract_text_from_pptx(uploaded_file)

            if raw_text.strip():
                st.session_state['summary'] = generate_ai_style_summary(
                    raw_text, 
                    summary_length, 
                    summary_format
                )
                st.session_state['file_name'] = uploaded_file.name
            else:
                st.warning("Tidak ditemukan teks yang dapat dibaca dari dokumen ini.")

if 'summary' in st.session_state:
    summary = st.session_state['summary']
    
    st.markdown("---")
    st.markdown(f"### {emojis[0]} Preview Lembar Rangkuman ({paper_style})")
    
    # Live Paper Preview
    st.markdown(f"""
        <div class="preview-paper">
            <h2 style="text-align: center; margin-top: 0;">{emojis[0]} Laaura's Resume {emojis[1]}</h2>
            <p style="text-align: center; font-weight: bold; margin-bottom: 25px; opacity: 0.8;">✨ Smart Document Summarizer ✨</p>
            <hr style="border: none; border-top: 1px dashed {theme['line_color']}; margin-bottom: 20px;">
            <div>{summary}</div>
        </div>
    """, unsafe_allow_html=True)
    
    # Preview Gambar
    if uploaded_images:
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown(f"#### {emojis[2]} Lampiran Gambar Tambahan ({len(uploaded_images)})")
        cols_img = st.columns(min(3, len(uploaded_images)))
        for idx, img_file in enumerate(uploaded_images):
            with cols_img[idx % 3]:
                st.image(img_file, caption=f"Lampiran {idx+1}", use_container_width=True)
                
    st.markdown("<br>", unsafe_allow_html=True)
    
    # Download Button
    pdf_bytes = create_custom_pdf(
        summary, 
        uploaded_images, 
        theme, 
        selected_font, 
        selected_size, 
        paper_style
    )
    st.download_button(
        label="📥 Download Version PDF (Custom Designed)",
        data=pdf_bytes,
        file_name=f"Laaura_Resume_{st.session_state['file_name']}.pdf",
        mime="application/pdf"
    )
